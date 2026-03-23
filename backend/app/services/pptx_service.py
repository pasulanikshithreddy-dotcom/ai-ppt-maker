from __future__ import annotations

import math
import re
from collections.abc import Mapping
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PptxPresentation
from pptx.slide import Slide
from pptx.util import Inches, Pt

from app.config.settings import Settings
from app.schemas.content_generation import (
    GeneratedPresentationContent,
    PresentationSlide,
)
from app.schemas.template import TemplateDefinition
from app.services.template_engine import TemplateEngine
from app.utils.runtime import new_id

ALIGNMENT_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}
FILENAME_SAFE_PATTERN = re.compile(r"[^a-zA-Z0-9]+")


class PptxTemplateNotFoundError(LookupError):
    """Raised when a requested presentation template does not exist."""


class PptxService:
    def __init__(self, settings: Settings) -> None:
        self.settings = settings
        self.template_engine = TemplateEngine(settings.template_catalog_path)

    @property
    def default_template_path(self) -> Path:
        return self.settings.ppt_template_dir / "default.pptx"

    def is_configured(self) -> bool:
        return True

    def readiness_detail(self) -> str:
        return (
            f"Template directory: {self.settings.ppt_template_dir}. "
            f"Output directory: {self.settings.generated_ppt_dir}."
        )

    def create_presentation(self) -> PptxPresentation:
        if self.default_template_path.exists():
            return Presentation(str(self.default_template_path))
        return Presentation()

    def get_template_definition(self, template_id: str) -> TemplateDefinition | None:
        return self.template_engine.get_template(template_id)

    def export_presentation(
        self,
        content: GeneratedPresentationContent | Mapping[str, object],
        *,
        template_id: str,
        file_stem: str | None = None,
    ) -> Path:
        presentation_content = GeneratedPresentationContent.model_validate(content)
        template = self.get_template_definition(template_id)
        if template is None:
            raise PptxTemplateNotFoundError(f"Template '{template_id}' was not found.")

        presentation = self.create_presentation()
        self._add_cover_slide(presentation, presentation_content.presentation_title, template)
        for slide_number, slide_content in enumerate(
            presentation_content.slides,
            start=1,
        ):
            self._add_content_slide(
                presentation,
                slide_content,
                template,
                slide_number=slide_number,
            )

        output_path = self._build_output_path(file_stem or presentation_content.presentation_title)
        presentation.save(str(output_path))
        return output_path.resolve()

    def _add_cover_slide(
        self,
        presentation: PptxPresentation,
        title: str,
        template: TemplateDefinition,
    ) -> None:
        slide = presentation.slides.add_slide(self._blank_layout(presentation))
        self._set_slide_background(slide, template.theme.secondary_color)
        self._apply_cover_style(slide, presentation, template)

        if template.layout.cover_style in {"bold_split", "hero_product"}:
            title_left = Inches(4.8)
            title_width = Inches(7.6)
        else:
            title_left = Inches(0.85)
            title_width = Inches(11.5)

        title_shape = slide.shapes.add_textbox(
            title_left,
            Inches(1.45),
            title_width,
            Inches(1.6),
        )
        title_frame = title_shape.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_frame.text = title
        self._style_paragraph(
            title_frame.paragraphs[0],
            font_family=template.theme.font_family,
            font_size=template.theme.title_font_size + 4,
            color_hex=template.theme.primary_color,
            alignment=template.layout.title_alignment,
            bold=True,
        )

        subtitle_shape = slide.shapes.add_textbox(
            title_left,
            Inches(3.15),
            title_width,
            Inches(0.75),
        )
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.text = f"{template.name} template"
        self._style_paragraph(
            subtitle_frame.paragraphs[0],
            font_family=template.theme.font_family,
            font_size=max(template.theme.body_font_size - 1, 12),
            color_hex=template.theme.primary_color,
            alignment=template.layout.title_alignment,
        )
        subtitle_frame.paragraphs[0].space_after = Pt(0)

    def _add_content_slide(
        self,
        presentation: PptxPresentation,
        slide_content: PresentationSlide,
        template: TemplateDefinition,
        *,
        slide_number: int,
    ) -> None:
        slide = presentation.slides.add_slide(self._blank_layout(presentation))
        self._set_slide_background(slide, template.theme.secondary_color)
        self._apply_content_decorations(slide, presentation, template)

        title_shape = slide.shapes.add_textbox(
            Inches(0.85),
            Inches(0.6),
            Inches(11.0),
            Inches(0.8),
        )
        title_frame = title_shape.text_frame
        title_frame.word_wrap = True
        title_frame.text = slide_content.title
        self._style_paragraph(
            title_frame.paragraphs[0],
            font_family=template.theme.font_family,
            font_size=template.theme.title_font_size,
            color_hex=template.theme.primary_color,
            alignment=template.layout.title_alignment,
            bold=True,
        )

        bullet_groups = self._split_bullets(
            slide_content.bullets,
            template.layout.content_columns,
        )
        self._render_bullet_groups(
            slide,
            presentation,
            bullet_groups,
            template,
        )

        self._add_speaker_notes(slide, slide_content.speaker_notes)
        if template.layout.show_page_numbers:
            self._add_page_number(
                slide,
                presentation,
                template,
                slide_number=slide_number,
            )

    @staticmethod
    def _blank_layout(presentation: PptxPresentation):
        if len(presentation.slide_layouts) > 6:
            return presentation.slide_layouts[6]
        return presentation.slide_layouts[-1]

    def _apply_cover_style(
        self,
        slide: Slide,
        presentation: PptxPresentation,
        template: TemplateDefinition,
    ) -> None:
        style = template.layout.cover_style
        if style in {"bold_split", "hero_product"}:
            self._add_rectangle(
                slide,
                left=0,
                top=0,
                width=int(presentation.slide_width * 0.31),
                height=presentation.slide_height,
                fill_hex=template.theme.primary_color,
            )
            self._add_rectangle(
                slide,
                left=int(presentation.slide_width * 0.31),
                top=Inches(6.85),
                width=int(presentation.slide_width * 0.69),
                height=Inches(0.25),
                fill_hex=template.theme.primary_color,
            )
            return

        if style in {"banner_top", "minimal_header"}:
            self._add_rectangle(
                slide,
                left=0,
                top=0,
                width=presentation.slide_width,
                height=Inches(0.85),
                fill_hex=template.theme.primary_color,
            )
            return

        if style in {"framed_title", "metric_panel", "editorial_grid"}:
            frame = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0.55),
                Inches(0.55),
                Inches(12.2),
                Inches(6.1),
            )
            frame.fill.background()
            frame.line.color.rgb = self._color(template.theme.primary_color)
            frame.line.width = Pt(1.75)
            self._add_rectangle(
                slide,
                left=Inches(0.55),
                top=Inches(5.95),
                width=Inches(12.2),
                height=Inches(0.24),
                fill_hex=template.theme.primary_color,
            )
            return

        self._add_rectangle(
            slide,
            left=Inches(0.85),
            top=Inches(0.8),
            width=Inches(1.75),
            height=Inches(0.16),
            fill_hex=template.theme.primary_color,
        )

    def _apply_content_decorations(
        self,
        slide: Slide,
        presentation: PptxPresentation,
        template: TemplateDefinition,
    ) -> None:
        if template.layout.use_accent_band:
            self._add_rectangle(
                slide,
                left=0,
                top=0,
                width=presentation.slide_width,
                height=Inches(0.18),
                fill_hex=template.theme.primary_color,
            )

        self._add_rectangle(
            slide,
            left=Inches(0.85),
            top=Inches(1.45),
            width=Inches(1.4),
            height=Inches(0.08),
            fill_hex=template.theme.primary_color,
        )

    def _render_bullet_groups(
        self,
        slide: Slide,
        presentation: PptxPresentation,
        bullet_groups: list[list[str]],
        template: TemplateDefinition,
    ) -> None:
        columns = len(bullet_groups)
        left_margin = Inches(0.9)
        right_margin = Inches(0.9)
        top = Inches(1.75)
        height = Inches(4.9)
        gutter = Inches(0.28)
        available_width = presentation.slide_width - left_margin - right_margin
        total_gutter = gutter * (columns - 1)
        column_width = int((available_width - total_gutter) / max(columns, 1))

        for column_index, bullet_group in enumerate(bullet_groups):
            left = left_margin + (column_index * (column_width + gutter))
            box = slide.shapes.add_textbox(left, top, column_width, height)
            frame = box.text_frame
            frame.clear()
            frame.word_wrap = True
            frame.vertical_anchor = MSO_ANCHOR.TOP
            self._populate_bullet_frame(frame, bullet_group, template)

    def _populate_bullet_frame(
        self,
        frame,
        bullets: list[str],
        template: TemplateDefinition,
    ) -> None:
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = f"\u2022 {bullet}"
            paragraph.space_after = Pt(10)
            paragraph.line_spacing = 1.15
            self._style_paragraph(
                paragraph,
                font_family=template.theme.font_family,
                font_size=template.theme.body_font_size,
                color_hex=template.theme.primary_color,
                alignment=template.layout.body_alignment,
            )

    def _add_speaker_notes(self, slide: Slide, speaker_notes: str) -> None:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.clear()
        notes_frame.text = speaker_notes

    def _add_page_number(
        self,
        slide: Slide,
        presentation: PptxPresentation,
        template: TemplateDefinition,
        *,
        slide_number: int,
    ) -> None:
        number_shape = slide.shapes.add_textbox(
            presentation.slide_width - Inches(1.1),
            presentation.slide_height - Inches(0.55),
            Inches(0.55),
            Inches(0.25),
        )
        number_frame = number_shape.text_frame
        number_frame.text = str(slide_number)
        self._style_paragraph(
            number_frame.paragraphs[0],
            font_family=template.theme.font_family,
            font_size=max(template.theme.body_font_size - 3, 10),
            color_hex=template.theme.primary_color,
            alignment="right",
            bold=True,
        )

    def _build_output_path(self, file_stem: str) -> Path:
        self.settings.generated_ppt_dir.mkdir(parents=True, exist_ok=True)
        safe_stem = self._sanitize_file_stem(file_stem)
        token = new_id("pptx").split("_", maxsplit=1)[-1][:10]
        return self.settings.generated_ppt_dir / f"{safe_stem}-{token}.pptx"

    @staticmethod
    def _sanitize_file_stem(value: str) -> str:
        cleaned = FILENAME_SAFE_PATTERN.sub("-", value.strip().lower()).strip("-")
        return cleaned[:60] or "presentation"

    @staticmethod
    def _split_bullets(bullets: list[str], column_count: int) -> list[list[str]]:
        actual_columns = max(1, min(column_count, len(bullets)))
        chunk_size = math.ceil(len(bullets) / actual_columns)
        return [
            bullets[index : index + chunk_size]
            for index in range(0, len(bullets), chunk_size)
        ]

    def _set_slide_background(self, slide: Slide, color_hex: str) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._color(color_hex)

    def _add_rectangle(
        self,
        slide: Slide,
        *,
        left: int,
        top: int,
        width: int,
        height: int,
        fill_hex: str,
    ) -> None:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left,
            top,
            width,
            height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._color(fill_hex)
        shape.line.fill.background()

    def _style_paragraph(
        self,
        paragraph,
        *,
        font_family: str,
        font_size: int,
        color_hex: str,
        alignment: str,
        bold: bool = False,
    ) -> None:
        paragraph.alignment = ALIGNMENT_MAP[alignment]
        for run in paragraph.runs:
            run.font.name = font_family
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = self._color(color_hex)

    @staticmethod
    def _color(value: str) -> RGBColor:
        return RGBColor.from_string(value)
