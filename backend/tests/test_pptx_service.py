from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from app.config.settings import Settings
from app.schemas.content_generation import (
    GeneratedPresentationContent,
    PresentationSlide,
)
from app.services.pptx_service import PptxService, PptxTemplateNotFoundError


def test_pptx_service_exports_structured_json_to_local_file(tmp_path: Path) -> None:
    settings = Settings(
        template_catalog_path=_template_catalog_path(),
        generated_ppt_dir=tmp_path,
    )
    service = PptxService(settings)
    content = GeneratedPresentationContent(
        presentation_title="AI PPT Maker Launch Plan",
        slides=[
            PresentationSlide(
                title="Problem",
                bullets=[
                    "Deck creation is still too manual",
                    "Teams lose time restructuring the same ideas",
                    "Visual consistency slips under deadline pressure",
                    "Speaker notes are often added too late",
                ],
                speaker_notes=(
                    "Frame the problem around speed, structure, and the friction "
                    "of manually polishing presentations."
                ),
            ),
            PresentationSlide(
                title="Solution",
                bullets=[
                    "Generate structured content from a topic",
                    "Apply reusable presentation templates",
                    "Export polished PowerPoint files quickly",
                ],
                speaker_notes=(
                    "Explain how AI PPT Maker shortens the path from raw idea "
                    "to a usable stakeholder deck."
                ),
            ),
            PresentationSlide(
                title="Outcome",
                bullets=[
                    "Faster first draft creation",
                    "Cleaner presentation storytelling",
                    "Consistent output across solo and team workflows",
                ],
                speaker_notes=(
                    "Close on user value: faster output, better polish, and more "
                    "repeatable deck quality."
                ),
            ),
        ],
    )

    output_path = service.export_presentation(
        content.model_dump(),
        template_id="startup_pitch",
    )

    assert output_path.exists()
    assert output_path.suffix == ".pptx"
    assert output_path.parent == tmp_path.resolve()

    presentation = Presentation(str(output_path))
    assert len(presentation.slides) == 4

    cover_title_shape = next(
        shape
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text") and shape.text == content.presentation_title
    )
    cover_run = cover_title_shape.text_frame.paragraphs[0].runs[0]
    assert cover_run.font.name == "Aptos"
    assert round(cover_run.font.size.pt) == 34

    notes_text = presentation.slides[1].notes_slide.notes_text_frame.text
    assert "speed, structure" in notes_text


def test_pptx_service_rejects_unknown_template(tmp_path: Path) -> None:
    settings = Settings(
        template_catalog_path=_template_catalog_path(),
        generated_ppt_dir=tmp_path,
    )
    service = PptxService(settings)
    content = GeneratedPresentationContent(
        presentation_title="Fallback Example",
        slides=[
            PresentationSlide(
                title="Slide One",
                bullets=["Point A", "Point B"],
                speaker_notes="Notes long enough to satisfy the presentation schema.",
            ),
            PresentationSlide(
                title="Slide Two",
                bullets=["Point C", "Point D"],
                speaker_notes="More speaker notes that validate successfully here.",
            ),
            PresentationSlide(
                title="Slide Three",
                bullets=["Point E", "Point F"],
                speaker_notes="Final notes block that keeps the payload valid.",
            ),
        ],
    )

    with pytest.raises(PptxTemplateNotFoundError):
        service.export_presentation(content, template_id="missing-template")


def _template_catalog_path() -> Path:
    return Path(__file__).resolve().parents[1] / "assets" / "templates" / "catalog.json"
